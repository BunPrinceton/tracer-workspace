#!/usr/bin/env python3
"""
Download all Google Docs/Sheets/Slides from a Drive folder as readable formats.

Converts:
- Google Docs → HTML (preserves formatting) or plain text
- Google Sheets → CSV
- Google Slides → plain text (extracts text from slides)
- Uploaded .docx/.pptx/.xlsx → Downloads as-is, then converts locally

Requirements:
    pip install google-auth google-auth-oauthlib google-api-python-client python-docx python-pptx openpyxl

Setup:
    1. Go to Google Cloud Console: https://console.cloud.google.com/
    2. Enable these APIs:
       - Google Drive API
       - Google Docs API (optional, for native docs)
    3. Create OAuth 2.0 credentials (Desktop app)
    4. Download as 'credentials.json' and place in this folder

Usage:
    python download_drive_docs.py --folder "FOLDER_ID" --output ./docs_output
    python download_drive_docs.py --folder "FOLDER_ID" --format html
    python download_drive_docs.py --folder "FOLDER_ID" --format txt
"""

import argparse
import os
import sys
import io
from pathlib import Path
from datetime import datetime

# Google API imports
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from googleapiclient.errors import HttpError

# Local conversion imports (for uploaded Office files)
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

# OAuth scopes - read-only access to Drive
SCOPES = [
    'https://www.googleapis.com/auth/drive.readonly',
]

# MIME type mappings
GOOGLE_MIME_TYPES = {
    'application/vnd.google-apps.document': {
        'name': 'Google Doc',
        'export_html': 'text/html',
        'export_txt': 'text/plain',
        'extension_html': '.html',
        'extension_txt': '.txt',
    },
    'application/vnd.google-apps.spreadsheet': {
        'name': 'Google Sheet',
        'export_html': 'text/csv',  # CSV is more useful than HTML for sheets
        'export_txt': 'text/csv',
        'extension_html': '.csv',
        'extension_txt': '.csv',
    },
    'application/vnd.google-apps.presentation': {
        'name': 'Google Slides',
        'export_html': 'text/plain',  # No good HTML export, use text
        'export_txt': 'text/plain',
        'extension_html': '.txt',
        'extension_txt': '.txt',
    },
}

OFFICE_MIME_TYPES = {
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': {
        'name': 'Word Doc (.docx)',
        'extension': '.docx',
        'converter': 'docx_to_text',
    },
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': {
        'name': 'PowerPoint (.pptx)',
        'extension': '.pptx',
        'converter': 'pptx_to_text',
    },
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': {
        'name': 'Excel (.xlsx)',
        'extension': '.xlsx',
        'converter': 'xlsx_to_text',
    },
    # Old Office formats
    'application/msword': {
        'name': 'Word Doc (.doc)',
        'extension': '.doc',
        'converter': None,  # Can't convert old .doc easily
    },
    'application/vnd.ms-powerpoint': {
        'name': 'PowerPoint (.ppt)',
        'extension': '.ppt',
        'converter': None,
    },
    'application/vnd.ms-excel': {
        'name': 'Excel (.xls)',
        'extension': '.xls',
        'converter': None,
    },
}

IMAGE_MIME_TYPES = {
    'image/png': '.png',
    'image/jpeg': '.jpg',
    'image/gif': '.gif',
    'image/webp': '.webp',
    'image/svg+xml': '.svg',
}


def get_credentials(creds_path: Path, token_path: Path):
    """Get or refresh OAuth credentials."""
    creds = None

    if token_path.exists():
        creds = Credentials.from_authorized_user_file(str(token_path), SCOPES)

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("  Refreshing expired token...")
            creds.refresh(Request())
        else:
            if not creds_path.exists():
                print(f"\n[ERR] ERROR: credentials.json not found at {creds_path}")
                print("\nTo set up Google Drive API:")
                print("1. Go to: https://console.cloud.google.com/")
                print("2. Create a project (or select existing)")
                print("3. Enable 'Google Drive API'")
                print("4. Go to Credentials → Create Credentials → OAuth 2.0 Client ID")
                print("5. Application type: Desktop app")
                print("6. Download JSON and save as 'credentials.json' here")
                sys.exit(1)

            print("\n[AUTH] Opening browser for Google authorization...")
            print("   (First time only - token will be saved for future runs)\n")
            flow = InstalledAppFlow.from_client_secrets_file(str(creds_path), SCOPES)
            creds = flow.run_local_server(port=0)

        # Save token for next run
        with open(token_path, 'w') as f:
            f.write(creds.to_json())
        print(f"  [OK] Token saved to {token_path}")

    return creds


def docx_to_text(file_path: Path) -> str:
    """Convert .docx to plain text."""
    if not HAS_DOCX:
        return f"[Cannot convert .docx - install python-docx: pip install python-docx]"

    doc = DocxDocument(str(file_path))
    paragraphs = [p.text for p in doc.paragraphs]
    return '\n\n'.join(paragraphs)


def pptx_to_text(file_path: Path) -> str:
    """Convert .pptx to plain text."""
    if not HAS_PPTX:
        return f"[Cannot convert .pptx - install python-pptx: pip install python-pptx]"

    prs = Presentation(str(file_path))
    text_parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_text = [f"=== Slide {i} ==="]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text)
        text_parts.append('\n'.join(slide_text))

    return '\n\n'.join(text_parts)


def xlsx_to_text(file_path: Path) -> str:
    """Convert .xlsx to plain text (CSV-like)."""
    if not HAS_XLSX:
        return f"[Cannot convert .xlsx - install openpyxl: pip install openpyxl]"

    wb = openpyxl.load_workbook(str(file_path), read_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = [f"=== Sheet: {sheet_name} ==="]

        for row in sheet.iter_rows(values_only=True):
            row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
            if row_text.strip():
                sheet_text.append(row_text)

        text_parts.append('\n'.join(sheet_text))

    return '\n\n'.join(text_parts)


def list_folder_contents(service, folder_id: str, path: str = "") -> list:
    """Recursively list all files in a Drive folder."""
    files = []
    page_token = None

    while True:
        query = f"'{folder_id}' in parents and trashed = false"
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='nextPageToken, files(id, name, mimeType, size, modifiedTime)',
            pageToken=page_token,
            pageSize=100
        ).execute()

        items = results.get('files', [])

        for item in items:
            item['path'] = path

            if item['mimeType'] == 'application/vnd.google-apps.folder':
                # Recursively get subfolder contents
                subfolder_path = f"{path}/{item['name']}" if path else item['name']
                print(f"  [DIR] Scanning folder: {subfolder_path}")
                files.extend(list_folder_contents(service, item['id'], subfolder_path))
            else:
                files.append(item)

        page_token = results.get('nextPageToken')
        if not page_token:
            break

    return files


def download_file(service, file_info: dict, output_dir: Path, format: str) -> dict:
    """Download a single file, converting if necessary."""
    file_id = file_info['id']
    file_name = file_info['name']
    mime_type = file_info['mimeType']
    file_path = file_info.get('path', '')

    result = {
        'name': file_name,
        'path': file_path,
        'mime_type': mime_type,
        'status': 'unknown',
        'output_file': None,
        'content': None,
    }

    # Create subfolder structure (sanitize path for Windows)
    if file_path:
        # Replace Windows-illegal characters in path
        safe_path = file_path.replace(':', '-').replace('?', '').replace('"', '').replace('<', '').replace('>', '').replace('|', '')
        output_subdir = output_dir / safe_path
        output_subdir.mkdir(parents=True, exist_ok=True)
    else:
        output_subdir = output_dir

    try:
        # Handle Google Workspace files (Docs, Sheets, Slides)
        if mime_type in GOOGLE_MIME_TYPES:
            type_info = GOOGLE_MIME_TYPES[mime_type]
            export_mime = type_info[f'export_{format}']
            extension = type_info[f'extension_{format}']

            # Clean filename and add extension
            safe_name = "".join(c for c in file_name if c.isalnum() or c in ' -_').strip()
            output_file = output_subdir / f"{safe_name}{extension}"

            request = service.files().export_media(fileId=file_id, mimeType=export_mime)
            content = request.execute()

            if isinstance(content, bytes):
                content = content.decode('utf-8', errors='replace')

            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)

            result['status'] = 'exported'
            result['output_file'] = str(output_file)
            result['content'] = content

        # Handle uploaded Office files
        elif mime_type in OFFICE_MIME_TYPES:
            type_info = OFFICE_MIME_TYPES[mime_type]
            extension = type_info['extension']
            converter = type_info['converter']

            # Download the binary file first
            safe_name = "".join(c for c in file_name if c.isalnum() or c in ' -_.').strip()
            if not safe_name.endswith(extension):
                safe_name += extension

            temp_file = output_subdir / safe_name

            request = service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)

            done = False
            while not done:
                status, done = downloader.next_chunk()

            with open(temp_file, 'wb') as f:
                f.write(fh.getvalue())

            # Convert to text if possible
            if converter:
                converter_func = globals()[converter]
                content = converter_func(temp_file)

                # Save text version
                text_file = output_subdir / f"{safe_name}.txt"
                with open(text_file, 'w', encoding='utf-8') as f:
                    f.write(content)

                result['status'] = 'converted'
                result['output_file'] = str(text_file)
                result['content'] = content
            else:
                result['status'] = 'downloaded_only'
                result['output_file'] = str(temp_file)
                result['content'] = f"[Binary file - no text converter available for {extension}]"

        # Handle images
        elif mime_type in IMAGE_MIME_TYPES:
            extension = IMAGE_MIME_TYPES[mime_type]
            safe_name = "".join(c for c in file_name if c.isalnum() or c in ' -_.').strip()
            if not any(safe_name.lower().endswith(ext) for ext in IMAGE_MIME_TYPES.values()):
                safe_name += extension

            output_file = output_subdir / safe_name

            request = service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)

            done = False
            while not done:
                status, done = downloader.next_chunk()

            with open(output_file, 'wb') as f:
                f.write(fh.getvalue())

            result['status'] = 'downloaded_image'
            result['output_file'] = str(output_file)
            result['content'] = f"[Image: {output_file}]"

        else:
            result['status'] = 'skipped'
            result['content'] = f"[Unsupported mime type: {mime_type}]"

    except HttpError as e:
        result['status'] = 'error'
        result['content'] = f"[Error: {e}]"
    except Exception as e:
        result['status'] = 'error'
        result['content'] = f"[Error: {type(e).__name__}: {e}]"

    return result


def main():
    parser = argparse.ArgumentParser(description='Download Google Drive folder contents as readable text')
    parser.add_argument('--folder', '-f', required=True, help='Google Drive folder ID')
    parser.add_argument('--output', '-o', default='./drive_docs_output', help='Output directory')
    parser.add_argument('--format', choices=['html', 'txt'], default='html', help='Export format for Google Docs')
    parser.add_argument('--credentials', '-c', default='credentials.json', help='Path to credentials.json')

    args = parser.parse_args()

    # Extract folder ID from URL if full URL provided
    folder_id = args.folder
    if 'drive.google.com' in folder_id:
        # Extract ID from URL like: https://drive.google.com/drive/folders/XXXXX?...
        if '/folders/' in folder_id:
            folder_id = folder_id.split('/folders/')[1].split('?')[0].split('&')[0]

    print("=" * 60)
    print("Google Drive Docs Downloader")
    print("=" * 60)
    print(f"Folder ID: {folder_id}")
    print(f"Output: {args.output}")
    print(f"Format: {args.format}")
    print("=" * 60 + "\n")

    # Setup paths
    script_dir = Path(__file__).parent
    creds_path = Path(args.credentials)
    if not creds_path.is_absolute():
        creds_path = script_dir / creds_path
    token_path = script_dir / 'drive_token.json'
    output_dir = Path(args.output)
    output_dir.mkdir(parents=True, exist_ok=True)

    # Authenticate
    print("Step 1: Authenticating...")
    creds = get_credentials(creds_path, token_path)
    service = build('drive', 'v3', credentials=creds)
    print("  [OK] Authenticated\n")

    # List folder contents
    print("Step 2: Scanning folder...")
    files = list_folder_contents(service, folder_id)
    print(f"  [OK] Found {len(files)} files\n")

    # Categorize files
    categories = {
        'google_docs': [],
        'office_files': [],
        'images': [],
        'other': [],
    }

    for f in files:
        if f['mimeType'] in GOOGLE_MIME_TYPES:
            categories['google_docs'].append(f)
        elif f['mimeType'] in OFFICE_MIME_TYPES:
            categories['office_files'].append(f)
        elif f['mimeType'] in IMAGE_MIME_TYPES:
            categories['images'].append(f)
        else:
            categories['other'].append(f)

    print("File breakdown:")
    print(f"  - Google Docs/Sheets/Slides: {len(categories['google_docs'])}")
    print(f"  - Office files (.docx/.pptx/.xlsx): {len(categories['office_files'])}")
    print(f"  - Images: {len(categories['images'])}")
    print(f"  - Other: {len(categories['other'])}")
    print()

    # Download all files
    print("Step 3: Downloading and converting files...")
    results = []
    total = len(files)

    for i, file_info in enumerate(files, 1):
        print(f"  [{i}/{total}] {file_info['name'][:50]}...", end=' ')
        result = download_file(service, file_info, output_dir, args.format)
        results.append(result)
        print(f"[{result['status']}]")

    print()

    # Summary
    print("=" * 60)
    print("SUMMARY")
    print("=" * 60)

    status_counts = {}
    for r in results:
        status_counts[r['status']] = status_counts.get(r['status'], 0) + 1

    for status, count in sorted(status_counts.items()):
        print(f"  {status}: {count}")

    print(f"\nOutput directory: {output_dir.absolute()}")
    print("=" * 60)

    # Create index file
    index_file = output_dir / '_index.txt'
    with open(index_file, 'w', encoding='utf-8') as f:
        f.write(f"Google Drive Download Index\n")
        f.write(f"Generated: {datetime.now().isoformat()}\n")
        f.write(f"Folder ID: {folder_id}\n")
        f.write(f"Total files: {len(results)}\n")
        f.write("=" * 60 + "\n\n")

        for r in results:
            f.write(f"File: {r['name']}\n")
            f.write(f"  Path: {r['path']}\n")
            f.write(f"  Type: {r['mime_type']}\n")
            f.write(f"  Status: {r['status']}\n")
            if r['output_file']:
                f.write(f"  Output: {r['output_file']}\n")
            f.write("\n")

    print(f"\n[OK] Index saved to: {index_file}")


if __name__ == '__main__':
    main()
